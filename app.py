import streamlit as st
import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

# --- 頁面基本設定 ---
st.set_page_config(
    page_title="流量異常鑑識儀表板",
    page_icon="🕵️‍♂️",
    layout="wide"
)

# --- 核心邏輯：資料清洗與計算 ---
@st.cache_data
def load_and_clean_data(file):
    try:
        df = pd.read_csv(file)
        numeric_cols = [
            'CTR（連結點閱率）', '曝光次數', '連結點擊次數', 
            '連結頁面瀏覽次數', 'CPM（每千次廣告曝光成本）', '花費金額 (TWD)'
        ]
        
        def clean_val(x):
            if isinstance(x, str):
                x = x.replace(',', '').replace('%', '')
                if x.strip() == '-': return 0
            return pd.to_numeric(x, errors='coerce')

        available_cols = [c for c in numeric_cols if c in df.columns]
        for col in available_cols:
            df[col] = df[col].apply(clean_val)
            
        if '連結頁面瀏覽次數' in df.columns and '連結點擊次數' in df.columns:
            df['LP_View_Rate'] = df.apply(
                lambda row: row['連結頁面瀏覽次數'] / row['連結點擊次數'] if row['連結點擊次數'] > 0 else 0, axis=1
            )
        else:
            df['LP_View_Rate'] = 0

        df_clean = df[df['曝光次數'] > 10].copy()
        return df_clean

    except Exception as e:
        st.error(f"檔案讀取錯誤: {e}")
        return None

# --- PPT 生成邏輯 (新功能) ---
def generate_pptx(ghost_df, flood_df, fig_ghost, fig_flood):
    prs = Presentation()

    # 1. 封面頁
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Meta 廣告流量異常鑑識報告"
    subtitle.text = "自動化偵測分析結果"

    # 輔助函數：新增表格到投影片
    def add_df_slide(prs, df, title_text):
        if df.empty: return
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank
        slide.shapes.title.text = title_text
        
        # 設定表格位置與大小
        rows, cols = min(df.shape[0] + 1, 11), df.shape[1] # 最多顯示 10 筆資料以免爆版
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 寫入標題列
        for i, col_name in enumerate(df.columns):
            cell = table.cell(0, i)
            cell.text = col_name
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True

        # 寫入數據 (取前 10 筆)
        for i in range(rows - 1):
            for j in range(cols):
                val = df.iloc[i, j]
                # 格式化數字
                if isinstance(val, float):
                    cell_text = f"{val:.2f}"
                else:
                    cell_text = str(val)
                
                cell = table.cell(i + 1, j)
                cell.text = cell_text
                cell.text_frame.paragraphs[0].font.size = Pt(9)

    # 輔助函數：新增圖表到投影片
    def add_chart_slide(prs, fig, title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank
        title = slide.shapes.title
        title.text = title_text
        
        # 將 Plotly 圖表轉為圖片串流 (需安裝 kaleido)
        img_bytes = fig.to_image(format="png", width=1000, height=600, scale=2)
        image_stream = BytesIO(img_bytes)
        
        # 貼上圖片
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), width=Inches(8))

    # 2. 幽靈點擊報告
    if not ghost_df.empty:
        # 整理表格欄位
        cols_export = ['天數', '廣告名稱', '連結點擊次數', '連結頁面瀏覽次數', 'CTR（連結點閱率）', 'LP_View_Rate']
        df_export = ghost_df[cols_export].head(10) # 只取前10筆
        add_df_slide(prs, df_export, "異常 A：幽靈點擊清單 (Top 10)")
        
        # 加圖表
        if fig_ghost:
            add_chart_slide(prs, fig_ghost, "異常 A：點擊品質診斷圖")

    # 3. 展示灌水報告
    if not flood_df.empty:
        cols_export = ['天數', '廣告名稱', '曝光次數', 'CTR（連結點閱率）', 'CPM（每千次廣告曝光成本）']
        df_export = flood_df[cols_export].head(10)
        add_df_slide(prs, df_export, "異常 B：展示灌水清單 (Top 10)")
        
        if fig_flood:
            add_chart_slide(prs, fig_flood, "異常 B：展示量 vs CTR 分布圖")

    # 存入記憶體
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# --- 側邊欄：參數控制 ---
st.sidebar.title("⚙️ 鑑識參數設定")

st.sidebar.subheader("1. 幽靈點擊偵測 (Ghost Clicks)")
threshold_ctr_high = st.sidebar.slider("CTR 異常高標 (%)", 2.0, 15.0, 4.0, 0.5)
threshold_quality_low = st.sidebar.slider("落地頁瀏覽率 低標 (Quality < X)", 0.1, 1.0, 0.5, 0.1)

st.sidebar.subheader("2. 展示灌水偵測 (Flooding)")
percentile_imp = st.sidebar.slider("高曝光定義 (PR值)", 50, 99, 75, 5)
threshold_ctr_low = st.sidebar.slider("CTR 異常低標 (%)", 0.1, 3.0, 1.5, 0.1)

# --- 主畫面 ---
st.title("🕵️‍♂️ 廣告流量異常鑑識系統")
st.markdown("上傳 CSV 報表，自動診斷流量異常，並支援 **一鍵匯出 PPT 報告**。")

uploaded_file = st.file_uploader("請上傳 CSV 報表檔案", type=['csv'])

if uploaded_file is not None:
    df = load_and_clean_data(uploaded_file)
    
    if df is not None:
        # --- 運算區 ---
        ghost_clicks = df[
            (df['CTR（連結點閱率）'] > threshold_ctr_high) & 
            (df['LP_View_Rate'] < threshold_quality_low)
        ].sort_values(by='CTR（連結點閱率）', ascending=False)

        imp_threshold_val = df['曝光次數'].quantile(percentile_imp / 100)
        flooding = df[
            (df['曝光次數'] > imp_threshold_val) & 
            (df['CTR（連結點閱率）'] < threshold_ctr_low)
        ].sort_values(by='曝光次數', ascending=False)

        # --- 繪圖區 (預先建立 fig 物件以便匯出) ---
        fig_ghost = None
        fig_flood = None

        if not df.empty:
            fig_ghost = px.scatter(
                df, x='連結點擊次數', y='連結頁面瀏覽次數', size='曝光次數', color='CTR（連結點閱率）',
                hover_data=['廣告名稱', '天數', 'LP_View_Rate'], title='點擊 vs. 到頁診斷', color_continuous_scale='Bluered'
            )
            max_val = df['連結點擊次數'].max()
            if pd.notnull(max_val):
                fig_ghost.add_shape(type="line", x0=0, y0=0, x1=max_val, y1=max_val, line=dict(color="Green", width=2, dash="dash"))

            fig_flood = px.scatter(
                df, x='曝光次數', y='CTR（連結點閱率）', size='CPM（每千次廣告曝光成本）', color='LP_View_Rate',
                hover_data=['廣告名稱', '天數', 'CPM（每千次廣告曝光成本）'], title='曝光 vs. CTR 診斷', color_continuous_scale='RdYlGn'
            )
            fig_flood.add_hline(y=threshold_ctr_low, line_dash="dash", line_color="red", annotation_text="低 CTR 警戒線")

        # --- 顯示區 ---
        st.markdown("---")
        col1, col2 = st.columns([1, 2])
        with col1:
            st.header("🚩 異常 A：幽靈點擊")
            st.metric("疑似異常數", f"{len(ghost_clicks)}")
        with col2:
            if fig_ghost: st.plotly_chart(fig_ghost, use_container_width=True)
        if not ghost_clicks.empty:
            st.dataframe(ghost_clicks[['天數', '廣告名稱', '曝光次數', 'CTR（連結點閱率）', 'LP_View_Rate']].style.format({'CTR（連結點閱率）': '{:.2f}%', 'LP_View_Rate': '{:.2%}'}))

        st.markdown("---")
        col3, col4 = st.columns([1, 2])
        with col3:
            st.header("🚩 異常 B：展示灌水")
            st.metric("疑似灌水數", f"{len(flooding)}")
        with col4:
            if fig_flood: st.plotly_chart(fig_flood, use_container_width=True)
        if not flooding.empty:
            st.dataframe(flooding[['天數', '廣告名稱', '曝光次數', 'CTR（連結點閱率）', 'CPM（每千次廣告曝光成本）']].style.format({'CTR（連結點閱率）': '{:.2f}%', 'CPM（每千次廣告曝光成本）': '{:.2f}'}))

        # --- 匯出按鈕區 ---
        st.markdown("---")
        st.header("📥 匯出報告")
        st.write("點擊下方按鈕，將當前的分析結果（含圖表與數據）下載為 PPT。")
        
        # 檢查是否安裝了 kaleido (圖表轉圖片需要)
        try:
            import kaleido
            can_export_charts = True
        except ImportError:
            can_export_charts = False
            st.warning("⚠️ 尚未安裝 'kaleido' 套件，匯出的 PPT 將不包含圖表，僅有數據表格。請執行 `pip install kaleido` 以啟用圖表匯出。")

        if st.button('生成 PPT 分析報告'):
            with st.spinner('正在生成簡報中，請稍候...'):
                try:
                    ppt_file = generate_pptx(ghost_clicks, flooding, fig_ghost if can_export_charts else None, fig_flood if can_export_charts else None)
                    
                    st.download_button(
                        label="⬇️ 下載 PPTX 檔案",
                        data=ppt_file,
                        file_name="Meta廣告異常分析報告.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    st.success("報告生成完畢！請點擊上方按鈕下載。")
                except Exception as e:
                    st.error(f"生成失敗: {e}")
else:
    st.info("請上傳檔案以開始分析。")
